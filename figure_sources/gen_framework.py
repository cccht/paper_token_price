"""
Generate academic paper-quality framework diagram for:
"Token Service Cross-Period Pricing & Demand Response Game Framework"
Three-end (User / Platform-Broker / Provider-Supply) architecture with algorithm layer.

Output: 16:9 widescreen .pptx for paper submission.
"""
import math
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.oxml.ns import qn

# ── Constants ──────────────────────────────────────────────
W = Cm(33.867)
H = Cm(19.05)

PRUSSIAN   = RGBColor(0x1B, 0x3A, 0x5C)   # dark blue: main titles, borders
STEEL      = RGBColor(0x2E, 0x6B, 0x9A)   # medium blue: boxes, headers
SKY        = RGBColor(0x5B, 0x9B, 0xD5)   # lighter blue: sub-boxes
LIGHT_BG   = RGBColor(0xD6, 0xE4, 0xF0)   # very light blue-grey: panel backgrounds
TEAL       = RGBColor(0x00, 0x7C, 0x7C)   # teal accent: user side
ORANGE     = RGBColor(0xE0, 0x7A, 0x2B)   # orange accent: platform side
GREEN      = RGBColor(0x2D, 0x8B, 0x4E)   # green accent: provider side
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MID_GRAY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)

FONT_BODY = "Calibri"
FONT_TITLE = "Calibri"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE


# ── Helper Functions ───────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None,
             border_color=None, border_width=Pt(1), corner_radius=None):
    """Add a rounded or sharp rectangle."""
    if corner_radius:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        # Default rounded corners are usually fine
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text,
                 font_size=Pt(9), color=DARK_GRAY, bold=False,
                 alignment=PP_ALIGN.LEFT, font_name=FONT_BODY,
                 anchor=MSO_ANCHOR.TOP):
    """Add a text box with single-style text."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    txbox.word_wrap = True
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def add_rich_box(slide, left, top, width, height, lines,
                 fill_color=None, border_color=None, border_width=Pt(1)):
    """Add a rounded rect with rich multi-line text.
    lines: list of (text, font_size, color, bold, alignment) tuples.
    """
    shape = add_rect(slide, left, top, width, height,
                     fill_color=fill_color, border_color=border_color,
                     border_width=border_width, corner_radius=True)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # Clear default paragraph
    for i, (text, fs, color, bold, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = fs
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT_BODY
        p.alignment = align
        p.space_before = Pt(1)
        p.space_after = Pt(1)
    return shape


def add_connector(slide, start_shape, end_shape, color=MID_GRAY,
                  width=Pt(1.5), dashed=False):
    """Add a straight connector arrow between two shapes (center to center)."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, 1, 1, 1, 1)  # dummy coords
    conn.begin_x = start_shape.left + start_shape.width // 2
    conn.begin_y = start_shape.top + start_shape.height
    conn.end_x   = end_shape.left + end_shape.width // 2
    conn.end_y   = end_shape.top
    conn.line.color.rgb = color
    conn.line.width = width
    if dashed:
        conn.line.dash_style = 2  # dash
    # Add arrowhead at end
    conn.end_x = end_shape.left + end_shape.width // 2
    conn.end_y = end_shape.top
    return conn


def add_arrow_between(slide, x1, y1, x2, y2, color=MID_GRAY,
                      width=Pt(1.5), dashed=False, label="",
                      label_font=Pt(7), label_color=MID_GRAY):
    """Add a straight arrow with optional label above midpoint."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                       x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = width
    if dashed:
        conn.line.dash_style = 2
    # arrowhead via XML
    ln = conn._element.find(qn('a:ln'))
    if ln is None:
        # For connectors, the line properties are stored differently
        pass

    if label:
        mx, my = (x1 + x2) // 2, (y1 + y2) // 2
        add_text_box(slide, mx + Cm(0.15), my - Cm(0.3),
                     Cm(2.5), Cm(0.6), label,
                     font_size=label_font, color=label_color,
                     alignment=PP_ALIGN.LEFT)
    return conn


def add_connector_with_arrow(slide, x1, y1, x2, y2, color=MID_GRAY,
                              width=Pt(1.3), dashed=False):
    """Add a connector line with arrowhead via XML manipulation."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                       x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = width
    if dashed:
        conn.line.dash_style = 2
    # Add tail and head end styles via XML
    spPr = conn._element.find(qn('p:spPr'))
    if spPr is not None:
        ln = spPr.find(qn('a:ln'))
        if ln is not None:
            tail = ln.find(qn('a:tailEnd'))
            if tail is None:
                tail = ln.makeelement(qn('a:tailEnd'), {})
                ln.append(tail)
            tail.set('type', 'triangle')
            tail.set('w', 'med')
            tail.set('len', 'med')
            head = ln.find(qn('a:headEnd'))
            if head is None:
                head = ln.makeelement(qn('a:headEnd'), {})
                ln.append(head)
            head.set('type', 'triangle')
            head.set('w', 'med')
            head.set('len', 'med')
    return conn


def add_dashed_arrow_between(slide, x1, y1, x2, y2, color=MID_GRAY,
                              width=Pt(1.3), label="", label_font=Pt(7),
                              label_color=MID_GRAY):
    """Dashed arrow with arrowhead on both ends for bidirectional."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                       x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = width
    conn.line.dash_style = 2  # dash
    # Add arrowhead via XML
    spPr = conn._element.find(qn('p:spPr'))
    if spPr is not None:
        ln = spPr.find(qn('a:ln'))
        if ln is not None:
            tail = ln.find(qn('a:tailEnd'))
            if tail is None:
                tail = ln.makeelement(qn('a:tailEnd'), {})
                ln.append(tail)
            tail.set('type', 'triangle')
            tail.set('w', 'med')
            tail.set('len', 'med')

    if label:
        mx, my = (x1 + x2) // 2, (y1 + y2) // 2
        add_text_box(slide, mx + Cm(0.2), my - Cm(0.3),
                     Cm(2.5), Cm(0.6), label,
                     font_size=label_font, color=label_color,
                     alignment=PP_ALIGN.LEFT)
    return conn


# ── Layout Grid ────────────────────────────────────────────
MARGIN = Cm(0.6)
PAGE_W = W
PAGE_H = H

TITLE_H = Cm(1.0)
ALGO_H  = Cm(2.8)
MAIN_TOP = TITLE_H + Cm(0.3)
MAIN_H   = PAGE_H - MAIN_TOP - ALGO_H - Cm(0.4)

COL_W = (PAGE_W - 4 * MARGIN) / 3
COL1_X = MARGIN
COL2_X = MARGIN + COL_W + MARGIN
COL3_X = MARGIN + 2 * (COL_W + MARGIN)

# ── 1. Title ────────────────────────────────────────────────
add_text_box(slide, MARGIN, Cm(0.15), PAGE_W - 2*MARGIN, TITLE_H,
             "Token Service Cross-Period Pricing & Demand Response Game Framework",
             font_size=Pt(16), color=PRUSSIAN, bold=True,
             alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

# ── 2. Three Main Columns Background ───────────────────────
# Subtle background panels for each end
col_bg_colors = {
    'user':     RGBColor(0xE8, 0xF4, 0xF0),  # teal tint
    'platform': RGBColor(0xFD, 0xF0, 0xE4),  # orange tint
    'provider': RGBColor(0xEA, 0xF5, 0xEA),  # green tint
}

bg_user     = add_rect(slide, COL1_X, MAIN_TOP, COL_W, MAIN_H,
                       fill_color=col_bg_colors['user'],
                       border_color=RGBColor(0xCC, 0xDD, 0xD6),
                       corner_radius=True)
bg_platform = add_rect(slide, COL2_X, MAIN_TOP, COL_W, MAIN_H,
                       fill_color=col_bg_colors['platform'],
                       border_color=RGBColor(0xE8, 0xD8, 0xC0),
                       corner_radius=True)
bg_provider = add_rect(slide, COL3_X, MAIN_TOP, COL_W, MAIN_H,
                       fill_color=col_bg_colors['provider'],
                       border_color=RGBColor(0xCC, 0xDD, 0xCC),
                       corner_radius=True)

# ── 3. Column Headers ──────────────────────────────────────
header_h = Cm(0.9)
header_y = MAIN_TOP + Cm(0.2)

add_rich_box(slide, COL1_X + Cm(0.15), header_y, COL_W - Cm(0.3), header_h,
             [("User Side", Pt(12), WHITE, True, PP_ALIGN.CENTER)],
             fill_color=TEAL, border_color=TEAL)

add_rich_box(slide, COL2_X + Cm(0.15), header_y, COL_W - Cm(0.3), header_h,
             [("Platform / Broker Side", Pt(12), WHITE, True, PP_ALIGN.CENTER)],
             fill_color=ORANGE, border_color=ORANGE)

add_rich_box(slide, COL3_X + Cm(0.15), header_y, COL_W - Cm(0.3), header_h,
             [("Provider / Supply Side", Pt(12), WHITE, True, PP_ALIGN.CENTER)],
             fill_color=GREEN, border_color=GREEN)

# ── 4. User Side Content ───────────────────────────────────
ux, uy = COL1_X + Cm(0.2), header_y + header_h + Cm(0.3)
uw = COL_W - Cm(0.4)

# Rigid User box
rigid_h = Cm(2.8)
rigid_box = add_rich_box(slide, ux, uy, uw, rigid_h, [
    ("Rigid Users", Pt(10), WHITE, True, PP_ALIGN.CENTER),
    ("", Pt(4), WHITE, False, PP_ALIGN.CENTER),
    ("Demand: R_t(p_t) = R̄_t · σ(β(v_r − p_t))", Pt(7.5), WHITE, False, PP_ALIGN.CENTER),
    ("", Pt(3), WHITE, False, PP_ALIGN.CENTER),
    ("Price-inelastic within willingness-to-pay", Pt(7), RGBColor(0xDD, 0xF0, 0xE8), False, PP_ALIGN.CENTER),
    ("Reduces only when p_t → v_r", Pt(7), RGBColor(0xDD, 0xF0, 0xE8), False, PP_ALIGN.CENTER),
], fill_color=TEAL, border_color=TEAL)

# Elastic User box
elastic_y = uy + rigid_h + Cm(0.25)
elastic_h = Cm(3.8)
elastic_box = add_rich_box(slide, ux, elastic_y, uw, elastic_h, [
    ("Elastic Users", Pt(10), WHITE, True, PP_ALIGN.CENTER),
    ("", Pt(4), WHITE, False, PP_ALIGN.CENTER),
    ("Logit discrete choice model:", Pt(7.5), WHITE, False, PP_ALIGN.CENTER),
    ("F_t(p) = F̄ · g · m · s_t", Pt(7.5), WHITE, True, PP_ALIGN.CENTER),
    ("", Pt(2), WHITE, False, PP_ALIGN.CENTER),
    ("s_t = exp(z_t) / Σ exp(z_τ)", Pt(7), RGBColor(0xDD, 0xF0, 0xE8), False, PP_ALIGN.CENTER),
    ("z_t = −α(p̃_t − p̄) + ln(a_t + ε)", Pt(7), RGBColor(0xDD, 0xF0, 0xE8), False, PP_ALIGN.CENTER),
    ("p̃_t = p_t + c_s(1 − a_t)", Pt(7), RGBColor(0xDD, 0xF0, 0xE8), False, PP_ALIGN.CENTER),
    ("", Pt(3), WHITE, False, PP_ALIGN.CENTER),
    ("Time preference a_t  ·  Migration cost c_s", Pt(7), RGBColor(0xDD, 0xF0, 0xE8), False, PP_ALIGN.CENTER),
    ("QoS feedback variant: z_t^feedback", Pt(7), RGBColor(0xDD, 0xF0, 0xE8), False, PP_ALIGN.CENTER),
], fill_color=TEAL, border_color=TEAL)

# Aggregated demand
agg_y = elastic_y + elastic_h + Cm(0.25)
agg_h = Cm(1.0)
agg_box = add_rich_box(slide, ux, agg_y, uw, agg_h, [
    ("Aggregate Demand:  D_t = R_t + F_t", Pt(9), WHITE, True, PP_ALIGN.CENTER),
], fill_color=RGBColor(0x00, 0x5F, 0x5F), border_color=RGBColor(0x00, 0x5F, 0x5F))


# ── 5. Platform Side Content ───────────────────────────────
px = COL2_X + Cm(0.2)
py = header_y + header_h + Cm(0.3)
pw = COL_W - Cm(0.4)

# Price vector
price_h = Cm(1.8)
price_box = add_rich_box(slide, px, py, pw, price_h, [
    ("Price Vector Publication", Pt(10), WHITE, True, PP_ALIGN.CENTER),
    ("", Pt(3), WHITE, False, PP_ALIGN.CENTER),
    ("p = (p₁, …, p_T)  ∈ [p̲, p̄]^T", Pt(8), WHITE, False, PP_ALIGN.CENTER),
    ("Fixed Average Price:  Σp_t / T = p̄", Pt(7.5), RGBColor(0xFF, 0xE8, 0xD0), False, PP_ALIGN.CENTER),
], fill_color=ORANGE, border_color=ORANGE)

# Stackelberg Leader
stk_y = py + price_h + Cm(0.2)
stk_h = Cm(1.0)
stk_box = add_rich_box(slide, px, stk_y, pw, stk_h, [
    ("Stackelberg Leader — First Mover", Pt(9), WHITE, True, PP_ALIGN.CENTER),
], fill_color=RGBColor(0xC0, 0x5A, 0x0A), border_color=RGBColor(0xC0, 0x5A, 0x0A))

# Profit maximization
profit_y = stk_y + stk_h + Cm(0.2)
profit_h = Cm(3.0)
profit_box = add_rich_box(slide, px, profit_y, pw, profit_h, [
    ("Profit Maximization (constrained)", Pt(10), WHITE, True, PP_ALIGN.CENTER),
    ("", Pt(3), WHITE, False, PP_ALIGN.CENTER),
    ("Π(p) = h Σ p_t D_t q(u_t)", Pt(8), WHITE, False, PP_ALIGN.CENTER),
    ("      − h Σ w D_t", Pt(7.5), RGBColor(0xFF, 0xE8, 0xD0), False, PP_ALIGN.CENTER),
    ("      − h Σ c_r(R̄_t − R_t)", Pt(7.5), RGBColor(0xFF, 0xE8, 0xD0), False, PP_ALIGN.CENTER),
    ("      − h Σ c_q[1−q(u_t)] R_t", Pt(7.5), RGBColor(0xFF, 0xE8, 0xD0), False, PP_ALIGN.CENTER),
    ("", Pt(3), WHITE, False, PP_ALIGN.CENTER),
    ("Constraints:", Pt(8), WHITE, True, PP_ALIGN.LEFT),
    ("   Bounded mean projection", Pt(7.5), RGBColor(0xFF, 0xE8, 0xD0), False, PP_ALIGN.LEFT),
    ("   Optional bill protection P_B(ρ)", Pt(7.5), RGBColor(0xFF, 0xE8, 0xD0), False, PP_ALIGN.LEFT),
], fill_color=ORANGE, border_color=ORANGE)

# Policy comparison
cmp_y = profit_y + profit_h + Cm(0.2)
cmp_h = Cm(1.9)
cmp_box = add_rich_box(slide, px, cmp_y, pw, cmp_h, [
    ("Pricing Strategies Compared", Pt(9), WHITE, True, PP_ALIGN.CENTER),
    ("", Pt(2), WHITE, False, PP_ALIGN.CENTER),
    ("1. Uniform pricing: p_t ≡ p̄", Pt(7.5), WHITE, False, PP_ALIGN.LEFT),
    ("2. Myopic dynamic: κ = 0 (search)", Pt(7.5), WHITE, False, PP_ALIGN.LEFT),
    ("3. QoS-aware dynamic: true κ", Pt(7.5), WHITE, False, PP_ALIGN.LEFT),
], fill_color=RGBColor(0xA0, 0x4A, 0x00), border_color=RGBColor(0xA0, 0x4A, 0x00))


# ── 6. Provider Side Content ───────────────────────────────
sx = COL3_X + Cm(0.2)
sy = header_y + header_h + Cm(0.3)
sw = COL_W - Cm(0.4)

# GPU cluster
gpu_h = Cm(2.0)
gpu_box = add_rich_box(slide, sx, sy, sw, gpu_h, [
    ("GPU Inference Cluster", Pt(10), WHITE, True, PP_ALIGN.CENTER),
    ("", Pt(3), WHITE, False, PP_ALIGN.CENTER),
    ("Fixed capacity: G", Pt(8), WHITE, False, PP_ALIGN.CENTER),
    ("Backends: vLLM 0.22.0 / Ollama 0.24.0", Pt(7.5), RGBColor(0xDD, 0xF5, 0xDD), False, PP_ALIGN.CENTER),
    ("RTX 4090 · Qwen2.5-0.5B/3B-Instruct", Pt(7.5), RGBColor(0xDD, 0xF5, 0xDD), False, PP_ALIGN.CENTER),
], fill_color=GREEN, border_color=GREEN)

# Utilization
util_y = sy + gpu_h + Cm(0.2)
util_h = Cm(1.0)
util_box = add_rich_box(slide, sx, util_y, sw, util_h, [
    ("Utilization:  u_t = D_t / G", Pt(9), WHITE, True, PP_ALIGN.CENTER),
], fill_color=RGBColor(0x1A, 0x6E, 0x35), border_color=RGBColor(0x1A, 0x6E, 0x35))

# QoS proxy
qos_y = util_y + util_h + Cm(0.2)
qos_h = Cm(2.5)
qos_box = add_rich_box(slide, sx, qos_y, sw, qos_h, [
    ("QoS Degradation Proxy", Pt(10), WHITE, True, PP_ALIGN.CENTER),
    ("", Pt(3), WHITE, False, PP_ALIGN.CENTER),
    ("q(u_t) = 1          if u_t ≤ ū", Pt(8), WHITE, False, PP_ALIGN.CENTER),
    ("q(u_t) = exp[−κ(u_t−ū)²]  if u_t > ū", Pt(8), WHITE, False, PP_ALIGN.CENTER),
    ("", Pt(2), WHITE, False, PP_ALIGN.CENTER),
    ("Continuous; models effective", Pt(7), RGBColor(0xDD, 0xF5, 0xDD), False, PP_ALIGN.CENTER),
    ("completion rate under congestion", Pt(7), RGBColor(0xDD, 0xF5, 0xDD), False, PP_ALIGN.CENTER),
], fill_color=GREEN, border_color=GREEN)

# Empirical validation
meas_y = qos_y + qos_h + Cm(0.2)
meas_h = Cm(2.5)
meas_box = add_rich_box(slide, sx, meas_y, sw, meas_h, [
    ("Empirical Validation (vLLM/Ollama)", Pt(10), WHITE, True, PP_ALIGN.CENTER),
    ("", Pt(3), WHITE, False, PP_ALIGN.CENTER),
    ("Metrics: TTFT, TPOT, throughput,", Pt(7.5), WHITE, False, PP_ALIGN.CENTER),
    ("  SLA compliance (TTFT < 0.5s)", Pt(7.5), WHITE, False, PP_ALIGN.CENTER),
    ("", Pt(2), WHITE, False, PP_ALIGN.CENTER),
    ("Heuristic QoS mapping:", Pt(7.5), RGBColor(0xDD, 0xF5, 0xDD), False, PP_ALIGN.CENTER),
    ("  SLA rate → normalized capacity → κ", Pt(7.5), RGBColor(0xDD, 0xF5, 0xDD), False, PP_ALIGN.CENTER),
], fill_color=RGBColor(0x1A, 0x5E, 0x2D), border_color=RGBColor(0x1A, 0x5E, 0x2D))


# ── 7. Inter-column arrows (main flows) ────────────────────

# User → Platform: Demand
arrow_y_forward = agg_y + agg_h // 2
arrow_y_back    = stk_y + stk_h // 2

# Platform → User: price signal (dashed, back direction)
add_dashed_arrow_between(
    slide,
    COL1_X + COL_W + Cm(0.05), arrow_y_back,
    COL2_X - Cm(0.05), arrow_y_back,
    color=ORANGE, width=Pt(2.0),
    label="Price Signal p_t", label_font=Pt(7), label_color=ORANGE)

# User → Platform: demand (solid, forward)
add_connector_with_arrow(
    slide,
    COL1_X + COL_W + Cm(0.05), agg_y + agg_h // 2,
    COL2_X - Cm(0.05), agg_y + agg_h // 2,
    color=TEAL, width=Pt(2.0))

add_text_box(slide, COL1_X + COL_W + Cm(0.2), agg_y + agg_h // 2 + Cm(0.15),
             Cm(2.0), Cm(0.5),
             "Demand D_t",
             font_size=Pt(7), color=TEAL)

# Platform → Provider: load allocation
add_connector_with_arrow(
    slide,
    COL2_X + COL_W, sy + gpu_h // 2,
    COL3_X, sy + gpu_h // 2,
    color=ORANGE, width=Pt(2.0))

add_text_box(slide, COL2_X + COL_W + Cm(0.1), sy + gpu_h // 2 - Cm(0.55),
             Cm(2.0), Cm(0.5),
             "Load D_t / Scheduling",
             font_size=Pt(7), color=ORANGE)

# Provider → Platform: QoS feedback (dashed, back)
add_dashed_arrow_between(
    slide,
    COL2_X + COL_W, qos_y + Cm(0.4),
    COL3_X - Cm(0.05), qos_y + Cm(0.4),
    color=GREEN, width=Pt(2.0),
    label="QoS q(u_t) / u_t", label_font=Pt(7), label_color=GREEN)

# User ↔ Provider: congestion externality (curved dashed)
add_dashed_arrow_between(
    slide,
    COL1_X + COL_W + Cm(0.1), agg_y + agg_h,
    COL3_X - Cm(0.1), util_y,
    color=MID_GRAY, width=Pt(1.2),
    label="Congestion Externality", label_font=Pt(6.5), label_color=MID_GRAY)


# ── 8. Game Theory & Algorithm Layer ───────────────────────
algo_y = MAIN_TOP + MAIN_H + Cm(0.3)
algo_w = PAGE_W - 2 * MARGIN

algo_bg = add_rect(slide, MARGIN, algo_y, algo_w, ALGO_H - Cm(0.1),
                   fill_color=LIGHT_BG, border_color=STEEL,
                   border_width=Pt(1.5), corner_radius=True)

add_text_box(slide, MARGIN + Cm(0.3), algo_y + Cm(0.1), Cm(18), Cm(0.6),
             "Algorithm & Game-Theoretic Foundation",
             font_size=Pt(10), color=PRUSSIAN, bold=True)

# Game theory sub-box
gt_x = MARGIN + Cm(0.3)
gt_y = algo_y + Cm(0.8)
gt_w = (algo_w - Cm(1.5)) / 2
gt_h = Cm(1.5)

add_rich_box(slide, gt_x, gt_y, gt_w, gt_h, [
    ("Game Structure", Pt(9), PRUSSIAN, True, PP_ALIGN.LEFT),
    ("Stackelberg Leader–Follower: platform first publishes p,", Pt(7.5), DARK_GRAY, False, PP_ALIGN.LEFT),
    ("users respond via aggregate logit choice (R_t, F_t)", Pt(7.5), DARK_GRAY, False, PP_ALIGN.LEFT),
    ("Congestion externality: overloaded GPU → QoS ↓ → effective service ↓", Pt(7.5), DARK_GRAY, False, PP_ALIGN.LEFT),
    ("Welfare decomposition: profit + rigid surplus + flex surplus − carbon cost", Pt(7.5), DARK_GRAY, False, PP_ALIGN.LEFT),
], fill_color=WHITE, border_color=STEEL)

# Algorithm sub-box
al_x = MARGIN + Cm(0.3) + gt_w + Cm(0.9)
al_y = algo_y + Cm(0.8)
al_w = gt_w
al_h = Cm(1.5)

add_rich_box(slide, al_x, al_y, al_w, al_h, [
    ("Numerical Optimization", Pt(9), PRUSSIAN, True, PP_ALIGN.LEFT),
    ("Multi-start SLSQP with equality, bound, and bill-protection constraints", Pt(7.5), DARK_GRAY, False, PP_ALIGN.LEFT),
    ("Bounded mean projection: bisection on λ via KKT optimality", Pt(7.5), DARK_GRAY, False, PP_ALIGN.LEFT),
    ("Candidates ranked by profit; solver convergence + constraint residual tracked", Pt(7.5), DARK_GRAY, False, PP_ALIGN.LEFT),
    ("Extensibility verified: T = 8 to 64, 5–40 starts, residual < 10⁻¹²", Pt(7.5), DARK_GRAY, False, PP_ALIGN.LEFT),
], fill_color=WHITE, border_color=STEEL)

# ── 9. Bottom annotation: QoS-driven feedback loop highlight ──

feedback_text_y = algo_y + ALGO_H - Cm(0.3)
add_text_box(slide, MARGIN, feedback_text_y, algo_w, Cm(0.45),
             "▸ QoS-aware pricing outperforms myopic by up to 12.53%  ·  QoS-min improves from 0.6632 (uniform) to 0.9918 (QoS-aware)  ·  Average-price constraint isolates scheduling gains from overall price hikes",
             font_size=Pt(7), color=PRUSSIAN, bold=False,
             alignment=PP_ALIGN.CENTER)

# ── 10. Save ───────────────────────────────────────────────
output_path = Path(__file__).resolve().with_name("framework_diagram.pptx")
prs.save(str(output_path))
print(f"Saved to {output_path}")
